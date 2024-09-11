import chainlit as cl
from chainlit.input_widget import Select, Switch
from chainlit.types import ThreadDict
from litellm import completion
import os
from dotenv import load_dotenv
from PyPDF2 import PdfReader
import docx
from typing import Optional, Dict, List
from botrun_ask_folder.fast_api.chainlit_custom_data_layer import FirestoreGCSDataLayer
import chainlit.data as cl_data
import pandas as pd
from pptx import Presentation
from botrun_ask_folder.split_txts import extract_text_from_pptx, convert_office_file
import mimetypes
from io import BytesIO
import base64
from openai import AsyncOpenAI
from chainlit.element import ElementBased

# 先不要使用 FirestoreGCSDataLayer
# cl_data._data_layer = FirestoreGCSDataLayer()
load_dotenv()

# 設置OpenAI客戶端
client = AsyncOpenAI()

# 設置全局變量
MAX_HISTORY_FOR_LLM = 10  # 設置發送給 LLM 的最大歷史記錄數量

MODEL_CONFIGS = {
    "openai/gpt-4o-mini": {
        "model": "openai/gpt-4o-mini",
        "provider": "OpenAI",
        "max_tokens": 4096,  # 調整為 gpt-4o-mini 的最大 token 數
    },
    "gemini/gemini-1.5-pro": {
        "model": "gemini/gemini-1.5-pro",
        "provider": "Gemini",
        "max_tokens": 4096,
    },
    "anthropic/claude-3-sonnet-20240229": {
        "model": "anthropic/claude-3-sonnet-20240229",
        "provider": "Anthropic",
        "max_tokens": 4096,
    },
}


@cl.oauth_callback
def oauth_callback(
    provider_id: str,
    token: str,
    raw_user_data: Dict[str, str],
    default_user: cl.User,
) -> Optional[cl.User]:
    return default_user


@cl.on_chat_start
async def start():
    user = cl.user_session.get("user")
    if not user:
        await cl.Message(content="認證失敗。請先登錄").send()
        return

    cl.user_session.set("chat_history", [])
    cl.user_session.set("current_model", "openai/gpt-4o-mini")  # 默認使用 gpt-4o-mini
    cl.user_session.set("file_contents", {})
    cl.user_session.set("audio_reply_enabled", False)  # 預設不啟用語音回覆

    settings = await cl.ChatSettings(
        [
            Select(
                id="model_selector",
                label="選擇模型",
                values=list(MODEL_CONFIGS.keys()),
                initial_index=list(MODEL_CONFIGS.keys()).index("openai/gpt-4o-mini"),
            ),
            Switch(id="audio_reply_enabled", label="啟用語音回覆", initial=False),
        ]
    ).send()

    welcome_message = f"歡迎，{user.display_name}！"
    await cl.Message(
        content=f"{welcome_message}您可以上傳檔案、直接開始對話，或按「麥克風」按鈕使用語音輸入。"
    ).send()


@cl.on_chat_resume
async def on_chat_resume(thread: ThreadDict):
    print(f"[chainlit_app][on_chat_resume] Resuming chat: {thread['id']}")
    cl.user_session.set("chat_history", thread.get("steps", []))
    cl.user_session.set("file_contents", {})  # Reset file contents

    # Restore file contents if any
    for element in thread.get("elements", []):
        if element["type"] == "file":
            await handle_file_upload(
                cl.File(name=element["name"], path=element["content"])
            )

    await cl.Message(content="聊天已恢復。您可以繼續之前的對話。").send()


@cl.on_settings_update
async def setup_agent(settings):
    cl.user_session.set("current_model", settings["model_selector"])
    cl.user_session.set("audio_reply_enabled", settings["audio_reply_enabled"])

    if settings["audio_reply_enabled"]:
        await cl.Message(
            content="語音回覆功能已啟用。系統將以語音形式回覆您的訊息。"
        ).send()
    else:
        await cl.Message(
            content="語音回覆功能已停用。系統將以文字形式回覆您的訊息。"
        ).send()


@cl.on_audio_chunk
async def on_audio_chunk(chunk: cl.AudioChunk):
    if chunk.isStart:
        buffer = BytesIO()
        buffer.name = f"input_audio.{chunk.mimeType.split('/')[-1]}"
        cl.user_session.set("audio_buffer", buffer)
        cl.user_session.set("audio_mime_type", chunk.mimeType)

    cl.user_session.get("audio_buffer").write(chunk.data)


@cl.on_audio_end
async def on_audio_end(elements: List[ElementBased]):
    audio_buffer: BytesIO = cl.user_session.get("audio_buffer")
    audio_buffer.seek(0)
    audio_file = audio_buffer.read()
    audio_mime_type: str = cl.user_session.get("audio_mime_type")

    input_audio_el = cl.Audio(
        mime=audio_mime_type, content=audio_file, name=audio_buffer.name
    )
    await cl.Message(
        author="您",
        type="user_message",
        content="",
        elements=[input_audio_el, *elements],
    ).send()

    whisper_input = (audio_buffer.name, audio_file, audio_mime_type)
    transcription = await speech_to_text(whisper_input)

    # 顯示轉錄文字
    await cl.Message(content=f"語音轉錄結果：{transcription}").send()

    # 處理轉錄後的文字
    await process_input_and_respond(transcription, elements)


@cl.step(type="tool")
async def speech_to_text(audio_file):
    response = await client.audio.transcriptions.create(
        model="whisper-1", file=audio_file
    )
    return response.text


async def text_to_speech(text: str, model: str):
    response = await client.audio.speech.create(model="tts-1", voice="nova", input=text)
    return "語音訊息", response.content


async def process_input_and_respond(
    input_content: str, elements: List[ElementBased] = None
):
    messages_for_llm, current_model, model_config = await process_input_and_files(
        input_content, elements
    )

    try:
        if any(isinstance(msg["content"], list) for msg in messages_for_llm):
            # 使用 OpenAI 的 gpt-4-vision-preview 模型
            response = completion(
                model="openai/gpt-4-vision-preview",
                api_key=os.getenv("OPENAI_API_KEY"),
                messages=messages_for_llm,
                max_tokens=model_config["max_tokens"],
                stream=True,
            )
        else:
            response = completion(
                model=model_config["model"],
                messages=messages_for_llm,
                api_key=os.getenv(f"{current_model.upper()}_API_KEY"),
                stream=True,
                max_tokens=model_config["max_tokens"],
            )

        answer_message = await cl.Message(content="").send()

        full_response = ""
        for chunk in response:
            if chunk.choices[0].delta.content is not None:
                content = chunk.choices[0].delta.content
                full_response += content
                await answer_message.stream_token(content)

        await answer_message.update()

        # 根據設置決定是否添加語音回覆
        if cl.user_session.get("audio_reply_enabled", False):
            output_name, output_audio = await text_to_speech(
                full_response, current_model
            )
            output_audio_el = cl.Audio(
                name=output_name,
                auto_play=True,
                mime="audio/mp3",
                content=output_audio,
            )
            answer_message.elements = [output_audio_el]
            await answer_message.update()

        chat_history = cl.user_session.get("chat_history", [])
        chat_history.append({"role": "user", "content": input_content})
        chat_history.append({"role": "assistant", "content": full_response})
        cl.user_session.set("chat_history", chat_history)

        return full_response

    except Exception as e:
        import traceback

        traceback.print_exc()
        error_message = (
            f"🙇‍♂️ 處理訊息時發生錯誤，請稍後再試或聯絡支援團隊。錯誤詳情：{str(e)}"
        )
        await cl.Message(content=error_message).send()
        return None


@cl.on_message
async def main(message: cl.Message):
    print(f"[chainlit_app][on_message] message: {message}")
    await process_input_and_respond(message.content, message.elements)


def encode_image(image_path):
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode("utf-8")


def handle_image_encode(element: cl.Image):
    images_content = cl.user_session.get("image_content", [])
    images_content.append(
        {
            "type": "image_url",
            "image_url": {
                "url": f"data:{element.mime};base64,{encode_image(element.path)}"
            },
        }
    )
    cl.user_session.set("image_content", images_content)


async def handle_file_upload(file: cl.File):
    content = ""
    # 獲取實際的 MIME 類型
    mime_type, _ = mimetypes.guess_type(file.path)

    if mime_type == "application/pdf" or file.mime == "application/pdf":
        content = extract_text_from_pdf(file.path)
    elif mime_type in [
        "application/msword",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ] or file.mime in [
        "application/msword",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ]:
        # 如果是 .doc 文件，先轉換為 .docx
        if file.name.lower().endswith(".doc"):
            converted_file_path = convert_office_file(file.path, ".docx")
            content = extract_text_from_docx(converted_file_path)
            os.remove(converted_file_path)
        else:
            content = extract_text_from_docx(file.path)
    elif mime_type in [
        "application/vnd.ms-excel",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ] or file.mime in [
        "application/vnd.ms-excel",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ]:
        df = pd.read_excel(file.path)
        content = df.to_string(index=False)
    elif mime_type in [
        "application/vnd.ms-powerpoint",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.oasis.opendocument.presentation",
    ] or file.mime in [
        "application/vnd.ms-powerpoint",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.oasis.opendocument.presentation",
    ]:
        try:
            # 總是先轉換為 .pptx 格式
            converted_file_path = convert_office_file(file.path, ".pptx")
            prs = Presentation(converted_file_path)
            content = ""
            for slide in prs.slides:
                content += extract_text_from_pptx(slide) + "\n\n"
            os.remove(converted_file_path)
        except Exception as e:
            print(f"Error processing PowerPoint file: {e}")
            await cl.Message(content=f"處理 PowerPoint 文件時出錯: {str(e)}").send()
            return
    elif (
        mime_type == "application/vnd.oasis.opendocument.spreadsheet"
        or file.mime == "application/vnd.oasis.opendocument.spreadsheet"
    ):
        converted_file_path = convert_office_file(file.path, ".xlsx")
        df = pd.read_excel(converted_file_path)
        content = df.to_string(index=False)
        os.remove(converted_file_path)
    elif (
        mime_type == "application/vnd.oasis.opendocument.presentation"
        or file.mime == "application/vnd.oasis.opendocument.presentation"
    ):
        converted_file_path = convert_office_file(file.path, ".pptx")
        prs = Presentation(converted_file_path)
        content = extract_text_from_pptx(prs)
        os.remove(converted_file_path)
    elif mime_type in [
        "application/rtf",
        "application/vnd.oasis.opendocument.text",
        "text/rtf",
    ] or file.mime in [
        "application/rtf",
        "application/vnd.oasis.opendocument.text",
        "text/rtf",
    ]:
        converted_file_path = convert_office_file(file.path, ".docx")
        content = extract_text_from_docx(converted_file_path)
        os.remove(converted_file_path)
    elif mime_type == "text/plain" or file.mime == "text/plain":
        with open(file.path, "r", encoding="utf-8") as f:
            content = f.read()
    else:
        await cl.Message(content=f"不支持的文件類型: {mime_type or file.mime}").send()
        return

    if not content:
        await cl.Message(content=f"無法提取文件 '{file.name}' 的內容。").send()
        return

    file_contents = cl.user_session.get("file_contents", {})
    file_contents[file.name] = content
    cl.user_session.set("file_contents", file_contents)
    await cl.Message(content=f"文件 '{file.name}' 已成功上傳和處理。").send()


def extract_text_from_pdf(file_path):
    with open(file_path, "rb") as file:
        reader = PdfReader(file)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
    return text


def extract_text_from_docx(file_path):
    doc = docx.Document(file_path)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text


async def process_input_and_files(
    message_content: str, elements: List[ElementBased] = None
):
    chat_history = cl.user_session.get("chat_history", [])
    current_model = cl.user_session.get("current_model")
    model_config = MODEL_CONFIGS[current_model]

    messages_for_llm = []
    for entry in chat_history[-MAX_HISTORY_FOR_LLM:]:
        messages_for_llm.append({"content": entry["content"], "role": entry["role"]})

    upload_has_file = False

    if elements:
        for element in elements:
            if isinstance(element, cl.File):
                if not upload_has_file:
                    cl.user_session.set("file_contents", {})
                    cl.user_session.set("image_content", [])
                    upload_has_file = True
                await handle_file_upload(element)
            elif isinstance(element, cl.Image):
                if not upload_has_file:
                    cl.user_session.set("file_contents", {})
                    cl.user_session.set("image_content", [])
                    upload_has_file = True
                handle_image_encode(element)

    file_contents = cl.user_session.get("file_contents", {})
    images_content = cl.user_session.get("image_content", [])

    # 添加文件內容作為系統消息
    if file_contents:
        combined_content = "\n\n".join(
            [
                f"<文件> '{name}'</文件>\n<內容>{content}</內容>"
                for name, content in file_contents.items()
            ]
        )
        messages_for_llm.append(
            {
                "content": f"以下是上傳的文件內容：\n\n{combined_content}\n\n請根據這些信息回答用戶的問題。",
                "role": "system",
            }
        )

    # 構建用戶消息
    user_message_content = [{"type": "text", "text": message_content}]
    if images_content:
        user_message_content.extend(images_content)
        user_message_content.append(
            {"type": "text", "text": "請根據上述圖片和文字回答問題。"}
        )

    messages_for_llm.append(
        {
            "role": "user",
            "content": user_message_content if images_content else message_content,
        }
    )

    return messages_for_llm, current_model, model_config
